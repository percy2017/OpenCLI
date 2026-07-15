from pathlib import Path

from pptx import Presentation


def _shape_text(shape) -> tuple[list[str], list[list[str]]]:
    """Return (paragraph_texts, table_rows) from one shape. Grouped shapes
    are recursed into. Empty entries are dropped."""
    paras: list[str] = []
    tables: list[list[str]] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                paras.append(t)
    if shape.has_table:
        row_strs: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_strs.append(" | ".join(cells))
        if row_strs:
            tables.append(row_strs)
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub in shape.shapes:
            sub_paras, sub_tables = _shape_text(sub)
            paras.extend(sub_paras)
            tables.extend(sub_tables)
    return paras, tables


def load_pptx(path: Path):
    prs = Presentation(path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        all_paras: list[str] = []
        all_tables: list[list[str]] = []
        for shape in slide.shapes:
            paras, tables = _shape_text(shape)
            all_paras.extend(paras)
            all_tables.extend(tables)

        if all_paras:
            yield {
                "text": "\n".join(all_paras),
                "slide": slide_num,
                "kind": "slide",
            }
        for row_strs in all_tables:
            yield {
                "text": "\n".join(row_strs),
                "slide": slide_num,
                "kind": "table",
            }

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = "\n".join(
                p.text.strip() for p in notes_tf.paragraphs if p.text.strip()
            )
            if notes_text:
                yield {
                    "text": notes_text,
                    "slide": slide_num,
                    "kind": "notes",
                }