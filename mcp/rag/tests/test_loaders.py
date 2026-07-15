"""Loader-level tests: assert each office extractor captures the content
that real documents put in non-paragraph blocks (tables, notes, etc.)."""
import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


def _docx_with_table(path: Path) -> None:
    doc = Document()
    doc.add_paragraph("Header line")
    table = doc.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Item"
    table.rows[0].cells[1].text = "Precio"
    table.rows[1].cells[0].text = "Hosting"
    table.rows[1].cells[1].text = "900"
    table.rows[2].cells[0].text = "Soporte"
    table.rows[2].cells[1].text = "0"
    doc.add_paragraph("Footer line")
    doc.save(path)


def _pptx_with_table_and_notes(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Title"
    table_shape = slide.shapes.add_table(rows=2, cols=2, left=Inches(1), top=Inches(2),
                                          width=Inches(4), height=Inches(1))
    table_shape.table.rows[0].cells[0].text = "Plan"
    table_shape.table.rows[0].cells[1].text = "Costo"
    table_shape.table.rows[1].cells[0].text = "Pro"
    table_shape.table.rows[1].cells[1].text = "500"
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = "Note: precio negociable"
    prs.save(path)


def _xlsx_with_rows(path: Path, n_rows: int) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Precios"
    ws.append(["Item", "Precio"])
    for i in range(n_rows):
        ws.append([f"Item-{i}", 100 + i])
    wb.save(path)


class DocxLoaderTests(unittest.TestCase):
    def test_captures_tables_alongside_paragraphs(self):
        from rag_mcp.loaders.docx import load_docx
        with tempfile.TemporaryDirectory() as tmp:
            p = Path(tmp) / "t.docx"
            _docx_with_table(p)
            recs = list(load_docx(p))
        kinds = [r.get("kind") for r in recs]
        self.assertIn("table", kinds, f"no table record yielded: {kinds}")
        # The price '900' must be in at least one record.
        joined = "\n".join(r["text"] for r in recs)
        self.assertIn("900", joined)
        self.assertIn("Hosting", joined)
        # Header and footer paragraphs still present.
        self.assertIn("Header line", joined)
        self.assertIn("Footer line", joined)


class PptxLoaderTests(unittest.TestCase):
    def test_captures_tables_and_notes(self):
        from rag_mcp.loaders.pptx import load_pptx
        with tempfile.TemporaryDirectory() as tmp:
            p = Path(tmp) / "t.pptx"
            _pptx_with_table_and_notes(p)
            recs = list(load_pptx(p))
        kinds = [r.get("kind") for r in recs]
        self.assertIn("table", kinds)
        self.assertIn("notes", kinds)
        joined = "\n".join(r["text"] for r in recs)
        self.assertIn("500", joined)
        self.assertIn("precio negociable", joined)


class XlsxLoaderTests(unittest.TestCase):
    def test_emits_one_record_per_row_block(self):
        from rag_mcp.loaders.xlsx import load_xlsx
        with tempfile.TemporaryDirectory() as tmp:
            p = Path(tmp) / "t.xlsx"
            _xlsx_with_rows(p, 60)
            recs = list(load_xlsx(p))
        # 60 rows + 1 header → 25+25+10 = 3 blocks
        self.assertEqual(len(recs), 3)
        for r in recs:
            self.assertEqual(r["kind"], "sheet_block")
            self.assertEqual(r["sheet"], "Precios")
            self.assertIn("Item:", r["text"])  # header is repeated per block
            self.assertIn("Precio:", r["text"])

    def test_header_only_sheet_is_still_indexed(self):
        from rag_mcp.loaders.xlsx import load_xlsx
        with tempfile.TemporaryDirectory() as tmp:
            p = Path(tmp) / "t.xlsx"
            wb = Workbook()
            ws = wb.active
            ws.append(["col1", "col2"])
            wb.save(p)
            recs = list(load_xlsx(p))
        self.assertEqual(len(recs), 1)
        self.assertEqual(recs[0]["kind"], "sheet_header_only")
        self.assertIn("col1", recs[0]["text"])


if __name__ == "__main__":
    unittest.main()
